from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- YARDIMCI FONKSİYONLAR ---
    def add_slide(title, content_list):
        slide_layout = prs.slide_layouts[1]  # Başlık ve İçerik düzeni
        slide = prs.slides.add_slide(slide_layout)
        
        # Başlık
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # İçerik
        tf = slide.shapes.placeholders[1].text_frame
        tf.text = content_list[0]  # İlk madde
        
        for item in content_list[1:]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            
        return slide

    # [cite_start]--- SLAYT 1: KAPAK [cite: 102-111] ---
    slide_layout = prs.slide_layouts[0] # Başlık Slaydı
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Python ile Stajyer Yerleştirme Simülasyonu"
    slide.placeholders[1].text = "Ahmet Yılmaz, Yunus Emre Erkuş, Ahmet Güldaş\n2025-2026 Dönem Projesi\nBursa Teknik Üniversitesi"

    # [cite_start]--- SLAYT 2: PROBLEM VE MOTİVASYON [cite: 121, 139] ---
    add_slide(
        "Problem Tanımı: Kıt Kaynaklar",
        [
            "📉 Kısıtlı Kontenjan: 150 Öğrenci vs. 40 Firma.",
            "🧩 Manuel Atama Zorluğu: GNO ve tercihleri elle eşleştirmek hataya açıktır.",
            "⚖️ Hedef: Toplam 'Memnuniyet Skorunu' maksimize eden adil dağıtım.",
            "🎯 Kapsam: Deterministik ve Stokastik (Olasılıksal) yaklaşımların kıyaslanması."
        ]
    )

    # [cite_start]--- SLAYT 3: TEKNOLOJİ YIĞINI [cite: 16, 22, 212] ---
    add_slide(
        "Yazılım Mimarisi ve Teknolojiler",
        [
            "🐍 Python 3.9 & Pandas: Veri manipülasyonu ve simülasyon motoru.",
            "🖥️ PyQt5 Arayüzü: Modern, kullanıcı dostu masaüstü uygulaması.",
            "⚡ Multithreading (QThread): Arayüz donmadan arka planda hesaplama.",
            "📊 Dinamik Raporlama: Anlık durum analizi ve CSV çıktıları."
        ]
    )

    # [cite_start]--- SLAYT 4: GREEDY ALGORİTMASI [cite: 41, 142, 145] ---
    add_slide(
        "Algoritma 1: Greedy (Açgözlü) Yaklaşım",
        [
            "Nasıl Çalışır? Öğrencileri GNO'ya göre sıralar, en başarılıyı ilk tercihine yerleştirir.",
            "✅ Avantaj: Çok hızlı (0.04 sn) ve deterministik.",
            "❌ Dezavantaj: Yerel (Local) en iyiye odaklanır, bütünün iyiliğini kaçırabilir.",
            "📉 Sonuç: 'İdare eder' bir memnuniyet skoru üretir (Baseline)."
        ]
    )

    # [cite_start]--- SLAYT 5: SIMULATED ANNEALING [cite: 63, 194, 202] ---
    s5 = add_slide(
        "Algoritma 2: Simulated Annealing (Benzetim Tavlaması)",
        [
            "Problem: Greedy ve Hill Climbing yerel tuzaklara (Local Optima) takılır.",
            "Çözüm: Metalurjiden esinlenen stokastik yaklaşım.",
            "🔥 Strateji: Başlangıçta (Yüksek Sıcaklık) kötü hamleleri kabul et.",
            "🧊 Sonuç: Zamanla soğuyarak Global Optimum noktasına ulaşır.",
            "📈 Skor Artışı: Greedy'e göre çok daha yüksek memnuniyet."
        ]
    )
    # Formülü not olarak ekleyelim
    notes_slide = s5.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Formül: P = exp(-DeltaE / T). Kötü hamleleri kabul etme olasılığı sıcaklığa bağlıdır."

    # [cite_start]--- SLAYT 6: DİNAMİK SİMÜLASYON [cite: 384, 385, 388] ---
    add_slide(
        "İnovasyon: Dinamik Red Simülasyonu",
        [
            "🎲 Kaos Testi: Gerçek hayatta her atama kabul edilmez.",
            "🔄 Döngü: Firmalar %10-%15 olasılıkla öğrenciyi reddeder (Mülakat Simülasyonu).",
            "⚙️ İyileşme: Reddedilenler havuza döner, sistem boşlukları tekrar doldurur.",
            "🛡️ Dayanıklılık: Sistem belirsizlik altında bile dengeye (equilibrium) ulaşır."
        ]
    )

    # [cite_start]--- SLAYT 7: SONUÇ TABLOSU [cite: 437, 438] ---
    slide_layout = prs.slide_layouts[5] # Sadece Başlık
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Karşılaştırmalı Analiz Sonuçları"
    
    # Tablo Ekleme
    rows, cols = 4, 4
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Başlıklar
    headers = ["Kriter", "Greedy", "Hill Climbing", "Simulated Annealing"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        
    # [cite_start]Veriler 
    data = [
        ["Çalışma Süresi", "~0.040 sn", "~6 sn", "~11 sn"],
        ["Memnuniyet Skoru", "12.260 (Baz)", "12.840", "13.120 (En İyi)"],
        ["İterasyon", "1 (Tek Geçiş)", "3000", "5000+"]
    ]
    
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, item in enumerate(row_data):
            table.cell(row_idx, col_idx).text = item

    # [cite_start]--- SLAYT 8: SONUÇ [cite: 467, 470] ---
    add_slide(
        "Sonuç ve Gelecek Çalışmalar",
        [
            "🏆 Başarı: Sezgisel yöntemler manuel atamaya göre üstün performans sağladı.",
            "🌍 Gerçekçilik: Dinamik red modülü ile sistemin dayanıklılığı kanıtlandı.",
            "🚀 Gelecek Planları:",
            "   - Coğrafi uzaklık kısıtlarının eklenmesi.",
            "   - Web tabanlı arayüze geçiş (Streamlit/Django)."
        ]
    )

    prs.save('Stajyer_Simulasyonu_Sunum.pptx')
    print("Sunum başarıyla oluşturuldu: Stajyer_Simulasyonu_Sunum.pptx")

if __name__ == "__main__":
    create_presentation()